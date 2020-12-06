from flask import Flask, render_template, redirect, request, session, make_response, send_file
from pptx import Presentation
from pptx.util import Pt
import wikipedia
from bs4 import BeautifulSoup
import requests
import re

app = Flask(__name__)

app.secret_key = b'testing'

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/download/<filename>")
def download(filename):
    return send_file(f'output/{filename}', filename, as_attachment=True)

@app.route("/generate/")
def generate():
    try:
        query = request.args.get('query')
        print(query)
        name = wikipedia.search(query, results=1)[0]
        url = "https://en.wikipedia.org/wiki/" + name.replace(" ", "_")
        summary = wikipedia.summary(name, sentences=3)
        print(url)
        # slide 1

        prs = Presentation('ppt/1.pptx')
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = name
        subtitle.text = "Auto-Generated Lesson Plan!"

        # slide 2

        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Summary of ' + name

        tf = body_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = summary
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(16)

        # slide n

        wiki = requests.get(url)
        wikisoup = BeautifulSoup(wiki.text, "lxml")
        toc = wikisoup.findAll("div", {"id":"toc" })[0]
        links = toc.findAll('a')
        temp = []
        sections = []
        for a in links:
            if any(a.text.split()[0][0] in s for s in temp):
                pass
            else:
                temp.append(a.text.split()[0])
                sections.append(a['href'][1:])

        for x in sections:
            try:
                title = wikisoup.find("span", {"id":x}).text
                text = re.sub("[\[].*?[\]]", "", wikisoup.find("span", {"id":x}).find_parent("h2").find_next_sibling("p").text)

                if text.strip() != "":
                    bullet_slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = slide.shapes
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]

                    title_shape.text = x.replace("_", " ").title()

                    tf = body_shape.text_frame
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = text.strip()
                    font = run.font
                    font.name = 'Calibri'
                    font.size = Pt(16)
            except:
                pass

        # last slide

        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Sources'

        tf = body_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = url
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(16)

        filename = name.replace(" ", "_").title() + '.pptx'
        prs.save("output/" + name.replace(" ", "_").title() + '.pptx')
    except:
        return render_template("error.html")

    return render_template("success.html", filename=filename, downloadurl=f"https://lesson.aru.wtf/download/{filename}")